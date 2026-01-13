"""
Ethereal Fashion 3D Rendering System - Full Proposal Presentation Generator
Creates a comprehensive PowerPoint presentation for the project proposal
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import os

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ETHEREAL FASHION"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(100, 200, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Zero-Gravity 3D Rendering System\nCyberpunk Meets Haute Couture"
    for para in subtitle_frame.paragraphs:
        para.font.size = Pt(28)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

    # Add team members
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    team_frame = team_box.text_frame
    team_frame.text = "Team Members: Person 1 | Person 2 | Person 3"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(20)
    team_para.font.color.rgb = RGBColor(200, 200, 200)
    team_para.alignment = PP_ALIGN.CENTER

    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 30)

def create_intro_slide(prs):
    """Create project introduction slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    title = slide.shapes.title
    title.text = "Project Introduction"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 200, 255)

    # Content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Vision Statement"

    p = tf.add_paragraph()
    p.text = "Revolutionizing digital fashion visualization through hyper-realistic 3D rendering technology"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Project Overview"

    p = tf.add_paragraph()
    p.text = "Create an ultra-realistic 3D rendering system that merges cyberpunk aesthetics with haute couture fashion"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Features ethereal zero-gravity environments, morphing iridescent fabrics, and holographic accessories"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Targets fashion designers, digital artists, NFT creators, and luxury brands"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Key Differentiator"

    p = tf.add_paragraph()
    p.text = "First-of-its-kind platform combining real-time physics simulation with cinematic-quality rendering"
    p.level = 1

def add_person1_slides(prs):
    """Person 1 - Problem, Background & Competitors"""

    # Slide 1: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 100, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Current Challenges in Fashion Visualization"

    p = tf.add_paragraph()
    p.text = "Traditional fashion photography is expensive (£5,000-£50,000 per shoot)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Physical prototyping wastes materials and time (4-6 weeks per sample)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Limited creative freedom with real-world physics constraints"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Digital fashion for metaverse/NFTs lacks photorealistic quality"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Growing demand for sustainable, virtual-first fashion solutions"
    p.level = 1

    # Slide 2: Why This Problem Matters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why This Problem is Relevant"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 100, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Market Drivers"

    p = tf.add_paragraph()
    p.text = "Digital fashion market projected to reach $50B by 2030"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "83% of fashion brands investing in digital transformation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Sustainability: Fashion industry accounts for 10% of global carbon emissions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Metaverse adoption: 400M+ users in virtual worlds"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "NFT fashion sales exceeded $137M in 2023"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Remote work increasing demand for virtual fashion showcases"
    p.level = 1

    # Slide 3: Background Research
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Background & Research"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 100, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Technology Foundation"

    p = tf.add_paragraph()
    p.text = "Physics-based rendering (PBR) enables photorealistic materials"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Real-time ray tracing revolutionized gaming graphics (2018+)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cloth simulation algorithms mature (Marvelous Designer, Blender)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "GPU compute power enables 8K rendering in reasonable timeframes"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Fashion Industry Trends"

    p = tf.add_paragraph()
    p.text = "Major brands (Gucci, Balenciaga, Dolce & Gabbana) entering digital fashion"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Virtual fashion shows becoming mainstream (post-2020)"
    p.level = 1

    # Slide 4: Competitor Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Competitor Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 100, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Key Competitors"

    competitors = [
        ("CLO3D / Marvelous Designer", "Industry standard for fashion design", "Limited rendering quality"),
        ("The Fabricant", "Digital fashion house", "Not a software platform"),
        ("DressX", "Digital fashion marketplace", "Basic AR try-ons only"),
        ("Browzwear", "3D fashion design software", "Enterprise-only, expensive"),
        ("Adobe Substance 3D", "General 3D texturing", "Not fashion-specific")
    ]

    for comp, pro, con in competitors:
        p = tf.add_paragraph()
        p.text = f"{comp}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"Pro: {pro} | Con: {con}"
        p.level = 2
        p.font.size = Pt(14)

    # Slide 5: Competitive Differentiation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Our Competitive Advantage"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 100, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "What Makes Us Different"

    p = tf.add_paragraph()
    p.text = "Only platform combining real-time editing with cinematic-quality output"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Specialized in avant-garde, impossible physics (zero-gravity, morphing fabrics)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "AI-powered procedural generation for unique patterns and effects"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "8K resolution output as standard (competitors max at 4K)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Built-in animation system for dynamic fashion presentations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Open ecosystem with plugin support and API access"
    p.level = 1

def add_person2_slides(prs):
    """Person 2 - Solution Design & System Requirements"""

    # Slide 1: Proposed Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Proposed Solution"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Ethereal Fashion Rendering Platform"

    p = tf.add_paragraph()
    p.text = "Cloud-based 3D rendering platform for hyper-realistic fashion visualization"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Combines real-time preview with offline cinematic rendering"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Browser-based editor + GPU render farm for processing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Pre-built templates for common fashion scenarios"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Target delivery: 4K preview in seconds, 8K final in minutes"
    p.level = 1

    # Slide 2: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame

    features = [
        ("Zero-Gravity Physics Engine", "Realistic floating cloth and hair simulation"),
        ("Iridescent Fabric Shader", "Color-shifting materials with real-time preview"),
        ("Galaxy Morphing System", "Fabrics transform into animated nebulae and constellations"),
        ("Holographic Accessories", "Levitating jewelry with particle trails"),
        ("Aurora Lighting System", "Volumetric atmospheric effects with color animation"),
        ("8K Ray-Traced Rendering", "Photorealistic output with 4096+ samples"),
        ("AI Pattern Generator", "Procedural texture generation for unique designs"),
        ("Animation Timeline", "Create dynamic fashion presentations")
    ]

    for feature, desc in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)

    # Slide 3: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Three-Tier Cloud Architecture"

    p = tf.add_paragraph()
    p.text = "Frontend Layer (Client)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Web-based editor (React + Three.js for real-time preview)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Asset management interface and timeline animator"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Backend Layer (API)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "REST API (Node.js/Python FastAPI) for job management"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Asset storage (AWS S3/Azure Blob) and authentication system"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Rendering Layer (GPU Cluster)"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Blender/Cycles X render farm with GPU nodes (RTX 4090/A100)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Queue management (Redis/RabbitMQ) and auto-scaling"
    p.level = 1

    # Slide 4: Hardware Requirements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Hardware Requirements"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Client-Side (Minimum)"

    p = tf.add_paragraph()
    p.text = "CPU: Quad-core 2.5GHz+ (Intel i5/AMD Ryzen 5)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "GPU: GTX 1660 or better (for real-time preview)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "RAM: 16GB DDR4"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Storage: 512GB SSD, 50Mbps internet"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Server-Side (Render Farm)"

    p = tf.add_paragraph()
    p.text = "GPU Nodes: 10x servers with NVIDIA RTX 4090 (24GB VRAM)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "CPU: AMD Threadripper/EPYC for scene preparation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Storage: 50TB NVMe SSD for asset library and cache"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Network: 10Gbps internal, 1Gbps external"
    p.level = 1

    # Slide 5: Software Requirements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Software Requirements"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Core Technologies"

    p = tf.add_paragraph()
    p.text = "3D Engine: Blender 4.0+ with Cycles X (open-source, Python API)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Frontend: React 18, Three.js, WebGL 2.0, TypeScript"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Backend: Python FastAPI + Node.js, PostgreSQL database"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Render Queue: Redis/RabbitMQ for job distribution"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cloud: AWS/Azure (EC2, S3, Lambda functions)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Texturing: Substance Designer SDK for procedural materials"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Physics: Bullet Physics + custom cloth simulation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "AI: Stable Diffusion for pattern generation (optional)"
    p.level = 1

    # Slide 6: Project Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Scope"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 255, 150)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "In Scope (Phase 1)"

    p = tf.add_paragraph()
    p.text = "Core rendering engine with zero-gravity physics"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "3 pre-made model templates (male/female/androgynous)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "10 fabric shader presets (iridescent, holographic, etc.)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Basic lighting and particle systems"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Web-based editor with real-time preview (1080p)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "8K still image rendering (no animation yet)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Out of Scope (Future Phases)"

    p = tf.add_paragraph()
    p.text = "Custom model import/rigging tools"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Full animation rendering (video output)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Mobile app versions"
    p.level = 1

def add_person3_slides(prs):
    """Person 3 - Planning & Team Information"""

    # Slide 1: Project Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Timeline"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 200, 100)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "12-Week Development Schedule"

    milestones = [
        ("Week 1-2: Research & Planning", "Requirements gathering, technology evaluation, architecture design"),
        ("Week 3-4: Core Engine Development", "Blender integration, basic physics simulation, shader development"),
        ("Week 5-6: Frontend Development", "Web editor UI, Three.js preview, asset management system"),
        ("Week 7-8: Backend & Infrastructure", "API development, render farm setup, cloud deployment"),
        ("Week 9-10: Integration & Testing", "End-to-end testing, performance optimization, bug fixing"),
        ("Week 11: Beta Testing & Refinement", "User testing with fashion designers, feedback implementation"),
        ("Week 12: Documentation & Launch", "User guides, marketing materials, soft launch")
    ]

    for milestone, details in milestones:
        p = tf.add_paragraph()
        p.text = milestone
        p.level = 0
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = details
        p.level = 1
        p.font.size = Pt(14)

    # Slide 2: Project Deliverables
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Deliverables"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 200, 100)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Technical Deliverables"

    p = tf.add_paragraph()
    p.text = "Fully functional web-based rendering platform (MVP)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Blender-based render farm with API integration"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Library of 10+ shader presets and 3 model templates"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cloud infrastructure (AWS/Azure) with auto-scaling"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Documentation & Marketing"

    p = tf.add_paragraph()
    p.text = "Technical documentation (API reference, architecture guide)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "User manual with video tutorials"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Marketing website and demo reel"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Sample Renders"

    p = tf.add_paragraph()
    p.text = "10 showcase renders demonstrating system capabilities"
    p.level = 1

    # Slide 3: Team Roles
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Team Roles & Responsibilities"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 200, 100)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Person 1: Project Lead & Research Specialist"

    p = tf.add_paragraph()
    p.text = "Market research, competitor analysis, problem definition"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Overall project coordination and stakeholder communication"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Person 2: Technical Architect & Backend Developer"

    p = tf.add_paragraph()
    p.text = "System architecture design, backend API development"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Render farm setup, cloud infrastructure management"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Person 3: Frontend Developer & UI/UX Designer"

    p = tf.add_paragraph()
    p.text = "Web editor interface design and development"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Real-time preview system (Three.js integration)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Shared Responsibilities"

    p = tf.add_paragraph()
    p.text = "All: Testing, documentation, 3D shader development"
    p.level = 1

    # Slide 4: Skills & Knowledge Acquisition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Skills & Knowledge Acquisition"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 200, 100)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Technical Skills Development"

    skills = [
        ("Advanced 3D Rendering", "Blender Cycles X, PBR materials, ray tracing optimization"),
        ("Physics Simulation", "Cloth dynamics, particle systems, zero-gravity simulation"),
        ("Shader Programming", "GLSL/OSL for custom material effects"),
        ("Cloud Architecture", "AWS/Azure deployment, auto-scaling, distributed computing"),
        ("Real-time Graphics", "WebGL, Three.js, GPU optimization"),
        ("API Design", "RESTful services, job queuing, microservices"),
        ("Fashion Industry Knowledge", "Fabric properties, haute couture trends, digital fashion markets"),
        ("Performance Optimization", "8K rendering efficiency, GPU memory management")
    ]

    for skill, details in skills:
        p = tf.add_paragraph()
        p.text = f"{skill}: {details}"
        p.level = 1
        p.font.size = Pt(14)

    # Slide 5: Mentor & Support
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Mentor & Support Resources"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 200, 100)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Project Mentor"

    p = tf.add_paragraph()
    p.text = "Name: Dr. [Mentor Name] / Industry Professional"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Expertise: Computer Graphics, 3D Rendering, Cloud Computing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Role: Technical guidance, weekly progress reviews, industry connections"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Advisory Support"

    p = tf.add_paragraph()
    p.text = "Fashion industry consultant for design validation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cloud infrastructure specialist (AWS/Azure certified)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Learning Resources"

    p = tf.add_paragraph()
    p.text = "Blender Studio, Substance Academy, AWS training programs"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Fashion-tech conferences and online communities"
    p.level = 1

def create_references_slide(prs):
    """Create references slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "References"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 255)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Technical Resources"

    references = [
        "Blender Foundation. (2024). Cycles X Rendering Engine Documentation",
        "Substance by Adobe. (2024). Procedural Material Design Guide",
        "Pharr, M., et al. (2023). Physically Based Rendering: Theory to Implementation",
        "Market Research",
        "McKinsey & Company. (2023). State of Fashion Technology Report",
        "Morgan Stanley. (2024). Digital Fashion Market Analysis",
        "Competitor Analysis",
        "CLO Virtual Fashion, The Fabricant, DressX, Browzwear (product reviews)",
        "Industry Publications",
        "WWD, Vogue Business, Fashion Innovation Agency reports"
    ]

    for ref in references:
        p = tf.add_paragraph()
        if ":" not in ref and len(ref.split()) < 4:
            p.text = ref
            p.font.bold = True
            p.level = 0
        else:
            p.text = ref
            p.level = 1
            p.font.size = Pt(14)

def main():
    """Main function to create the presentation"""
    print("=" * 70)
    print("CREATING ETHEREAL FASHION PROPOSAL PRESENTATION")
    print("=" * 70)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n✓ Creating title slide...")
    create_title_slide(prs)

    print("✓ Creating introduction slide...")
    create_intro_slide(prs)

    print("✓ Adding Person 1 slides (Problem, Background & Competitors)...")
    add_person1_slides(prs)

    print("✓ Adding Person 2 slides (Solution Design & System Requirements)...")
    add_person2_slides(prs)

    print("✓ Adding Person 3 slides (Planning & Team Information)...")
    add_person3_slides(prs)

    print("✓ Creating references slide...")
    create_references_slide(prs)

    # Save presentation
    output_file = '/vercel/sandbox/Ethereal_Fashion_Proposal.pptx'
    prs.save(output_file)

    print(f"\n{'=' * 70}")
    print(f"✓ PRESENTATION CREATED SUCCESSFULLY!")
    print(f"{'=' * 70}")
    print(f"\nFile: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"\nSlide breakdown:")
    print(f"  • Title & Introduction: 2 slides")
    print(f"  • Person 1 (Problem/Background/Competitors): 5 slides")
    print(f"  • Person 2 (Solution/Architecture/Requirements): 6 slides")
    print(f"  • Person 3 (Planning/Timeline/Team): 5 slides")
    print(f"  • References: 1 slide")
    print(f"\n{'=' * 70}")

if __name__ == "__main__":
    main()
